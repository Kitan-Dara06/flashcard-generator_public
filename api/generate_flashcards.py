# api/generate_flashcards.py
import json
import io
import os
import re
import base64

from typing import List
import pdfplumber
from http.server import BaseHTTPRequestHandler
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate, HumanMessagePromptTemplate
from langchain_core.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
import logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)  
# ----- Pydantic models -----
class Flashcard(BaseModel):
    question: str = Field(description="The question on the flashcard")
    answer: str = Field(description="The answer on the flashcard")

class FlashcardList(BaseModel):
    flashcards: List[Flashcard] = Field(description="List of flashcards")


# ----- JSON + Parsing Helpers -----
def clean_json_output(text: str):
    """Fix common JSON issues like trailing commas before parsing."""
    return re.sub(r",(\s*[}\]])", r"\1", text)

def safe_parse_flashcards(result):
    """Ensure flashcards always have 'question' and 'answer' fields."""
    repaired = []
    flashcards = getattr(result, "flashcards", []) if hasattr(result, "flashcards") else result.get("flashcards", [])

    for c in flashcards:
        if isinstance(c, dict):
            q = c.get("question", "").strip()
            a = c.get("answer", "").strip() or "Answer not provided in text."
        else:  # Pydantic model case
            q = getattr(c, "question", "").strip()
            a = getattr(c, "answer", "").strip() or "Answer not provided in text."
        repaired.append({"question": q, "answer": a})
    return repaired

def parse_with_json_fallback(raw_output: str):
    """Fallback: force JSON.loads."""
    try:
        cleaned = clean_json_output(raw_output)
        data = json.loads(cleaned)
        return safe_parse_flashcards(data)
    except Exception as e:
        print(f"⚠️ JSON fallback failed: {e}")
        return []

def try_parse_flashcards(raw_output: str):
    """Try Pydantic parsing first, then fallback to cleaned JSON."""
    parser = PydanticOutputParser(pydantic_object=FlashcardList)
    try:
        cleaned = clean_json_output(raw_output)
        parsed = parser.parse(cleaned)
        if not getattr(parsed, "flashcards", None):
            raise ValueError("Parsed object missing flashcards")
        return safe_parse_flashcards(parsed)
    except Exception as e:
        print(f"⚠️ Parser failed, using JSON fallback: {e}")
        return parse_with_json_fallback(raw_output)





# ----- Text extraction -----
def extract_text_from_pdf(file_content, max_pages=100):
    try:
        with pdfplumber.open(io.BytesIO(file_content)) as pdf:
            total_pages = len(pdf.pages)
            if total_pages > max_pages:
                return {
                    "error": "document_too_long",
                    "message": f"Your document has {total_pages} pages. Max allowed: {max_pages}.",
                    "page_count": total_pages,
                    "max_allowed": max_pages
                }
            return "\n".join([page.extract_text() or "" for page in pdf.pages])
    except Exception as e:
        return {
            "error": "pdf_extraction_failed",
            "message": f"Error reading PDF: {str(e)}"
        }

def extract_text_from_docx(file_content):
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_content))
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        return {
            "error": "docx_extraction_failed",
            "message": f"Error reading DOCX: {str(e)}"
        }
def filter_valid_flashcards(flashcards):
    valid = []
    for item in flashcards:
        if isinstance(item, dict):
            q = item.get("question", "").strip()
            a = item.get("answer", "").strip()
            if q and a:
                valid.append({"question": q, "answer": a})
    return valid
    
def extract_text_from_pptx(file_content):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return {
            "error": "pptx_extraction_failed",
            "message": f"Error reading PPTX: {str(e)}"
        }


# ----- Flashcard generation -----

def generate_flashcards(text, api_key):
    prompt = ChatPromptTemplate.from_messages([
        HumanMessagePromptTemplate.from_template(
            """You are a flashcard generator for theory-based subjects.
Output a valid JSON object with a key "flashcards" containing a list of flashcards.
Each flashcard must have:
  - "question": a clear, concise question (string)
  - "answer": a 2–3 sentence explanatory answer (string)
Stay strictly factual, based only on the provided text.
If the text contains no usable information, output {"flashcards": []}.
Do not explain, apologize, or return any text outside the JSON object.

Text:
{input_text}"""
        )
    ])
    
    llm = ChatOpenAI(
        model="gpt-4o",
        temperature=0.3,
        api_key=api_key,
        response_format={"type": "json_object"},
        max_tokens=4000
    )
    chain = prompt | llm 

    CHARS_PER_TOKEN = 4  
    MAX_INPUT_TOKENS = 6000
    chunk_size = MAX_INPUT_TOKENS * CHARS_PER_TOKEN  
    
    chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    all_flashcards = [] 
    
    for i, chunk in enumerate(chunks, 1):
        if not chunk.strip():
            continue
            
        try:
            logger.info(f"Processing chunk {i}/{len(chunks)} ({len(chunk)} chars)")
            result = chain.invoke({"input_text": chunk})
            raw = result.content.strip()
            
            data = json.loads(raw)
            
            flashcards_raw = data.get("flashcards", [])
            if not filter_valid_flashcards(flashcards_raw):
                print("Invalid flashcard structure")
                continue
            flashcards = safe_parse_flashcards(flashcards_raw)
            if flashcards:  # Only add if we got valid flashcards
                all_flashcards.extend(flashcards)
                logger.info(f"Chunk {i}: Generated {len(flashcards)} flashcards")
            else:
                logger.warning(f"Chunk {i}: No valid flashcards generated")
        except Exception as e:
            logger.error(f"Error processing chunk {i}: {e}")
            continue  # Skip failed chunks

    return all_flashcards
           
# ----- Core handler logic -----
def lambda_handler(event):
    try:
        if event["httpMethod"] == "OPTIONS":
            return {
                "statusCode": 200,
                "headers": {
                    "Access-Control-Allow-Origin": "*",
                    "Access-Control-Allow-Methods": "POST, OPTIONS",
                    "Access-Control-Allow-Headers": "Content-Type",
                },
                "body": ""
            }

        if event["httpMethod"] != "POST":
            return {
                "statusCode": 405,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"success": False, "error": "Method not allowed"})
            }

        # Get OpenAI API key
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            return {
                "statusCode": 500,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"success": False, "error": "OpenAI API key not configured"})
            }

        # Normalize headers (case-insensitive)
        headers = {k.lower(): v for k, v in event["headers"].items()}
        content_type = headers.get("content-type", "")
        if "application/json" not in content_type:
            return {
                "statusCode": 400,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"success": False, "error": "Content-Type must be application/json"})
            }

        # Parse JSON body
        body = json.loads(event["body"])
        file_content = base64.b64decode(body["file_content"])
        file_type = body["file_type"]

        # Extract text
        if file_type == "application/pdf":
            text = extract_text_from_pdf(file_content)
        elif file_type == "text/plain":
            text = file_content.decode("utf-8")
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = extract_text_from_docx(file_content)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_pptx(file_content)
        else:
            return {
                "statusCode": 400,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"success": False, "error": "Unsupported file type"})
            }
        if isinstance(text, dict) and "error" in text:
            return {
            "statusCode": 400,
            "headers": { "Content-Type": "application/json",
                    "Access-Control-Allow-Origin": "*"},
            "body": json.dumps({ "success": False, **text })
            }

        if not text.strip():
            return {
                "statusCode": 400,
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({"success": False, "error": "Could not extract text from file"})
            }

        # Generate flashcards
        flashcards = generate_flashcards(text, api_key)
        return {
            "statusCode": 200,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*"
            },
            "body": json.dumps({
                "success": True,
                "flashcards": flashcards,
                "count": len(flashcards)
            })
        }

    except Exception as e:
        import traceback
        print("Unhandled error:", str(e))
        print(traceback.format_exc())
        return {
            "statusCode": 500,
            "headers": {
                "Content-Type": "application/json",
                "Access-Control-Allow-Origin": "*"
            },
            "body": json.dumps({
                "success": False,
                "error": str(e),
                "trace": traceback.format_exc()
            })
        }


# ----- Vercel entrypoint -----
class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        content_length = int(self.headers.get("Content-Length", 0))
        body = self.rfile.read(content_length).decode("utf-8")
        event = {
            "httpMethod": "POST",
            "headers": dict(self.headers),
            "body": body
        }
        response = lambda_handler(event)

        self.send_response(response["statusCode"])
        for k, v in response.get("headers", {}).items():
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(response["body"].encode())

    def do_OPTIONS(self):
        response = lambda_handler({"httpMethod": "OPTIONS", "headers": {}, "body": ""})
        self.send_response(response["statusCode"])
        for k, v in response.get("headers", {}).items():
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(response["body"].encode())

